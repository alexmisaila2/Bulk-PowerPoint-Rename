import os
import re
import unicodedata
import win32com.client  # Used to handle .ppt files with PowerPoint
from pptx import Presentation  # Used to handle .pptx files
import shutil

# Define log file path inside CantariBun folder
LOG_DIR = r"C:\Users\maste\Desktop\CantariBun"
LOG_FILE = os.path.join(LOG_DIR, "process_log.txt")

# Ensure log directory exists
if not os.path.exists(LOG_DIR):
    os.makedirs(LOG_DIR)

# Function to log messages
def log_message(message):
    with open(LOG_FILE, "a", encoding="utf-8") as log:
        log.write(message + "\n")
    print(message)

# Function to remove diacritics
def remove_diacritics(text):
    return ''.join((c if not unicodedata.category(c) == 'Mn' else '') for c in unicodedata.normalize('NFD', text))

# Function to remove special characters (excluding spaces and alphanumeric)
def remove_special_characters(text):
    return re.sub(r'[^\w\s\-]', '', text).strip()

# Function to extract first 6 words from a .pptx file
def get_first_6_words_from_pptx(file_path, slide_index):
    try:
        presentation = Presentation(file_path)
        slides = list(presentation.slides)
        if slide_index - 1 < len(slides):
            slide = slides[slide_index - 1]
            text = " ".join([shape.text for shape in slide.shapes if hasattr(shape, "text")])
            return " ".join(text.split()[:6])
    except Exception as e:
        log_message(f"ERROR reading .pptx file {file_path}: {e}")
    return ""

# Function to extract first 6 words from a .ppt file
def get_first_6_words_from_ppt(file_path, slide_index):
    try:
        powerpoint = win32com.client.Dispatch("PowerPoint.Application")
        presentation = powerpoint.Presentations.Open(file_path, WithWindow=False)
        
        if slide_index <= presentation.Slides.Count:
            slide = presentation.Slides(slide_index)
            text = " ".join([shape.TextFrame.TextRange.Text for shape in slide.Shapes if shape.HasTextFrame and shape.TextFrame.HasText])
        else:
            text = ""
        
        presentation.Close()
        powerpoint.Quit()
        return " ".join(text.split()[:6])
    except Exception as e:
        log_message(f"ERROR reading .ppt file {file_path}: {e}")
    return ""

# Function to extract text based on file type
def get_first_6_words(file_path, slide_index):
    if file_path.endswith(".pptx"):
        return get_first_6_words_from_pptx(file_path, slide_index)
    elif file_path.endswith(".ppt"):
        return get_first_6_words_from_ppt(file_path, slide_index)
    return ""

# Function to rename and copy files
def rename_and_copy_files(src_dir, dest_dir):
    for root, dirs, files in os.walk(src_dir):
        relative_path = os.path.relpath(root, src_dir)
        dest_folder = os.path.join(dest_dir, relative_path)
        os.makedirs(dest_folder, exist_ok=True)

        for file in files:
            if file.endswith(".ppt") or file.endswith(".pptx"):
                src_file_path = os.path.join(root, file)
                log_message(f"Processing file: {src_file_path}")

                first_6_words_slide_1 = remove_special_characters(remove_diacritics(get_first_6_words(src_file_path, 1)))
                first_6_words_slide_2 = remove_special_characters(remove_diacritics(get_first_6_words(src_file_path, 2)))

                new_file_name = f"{first_6_words_slide_1} ({first_6_words_slide_2}){os.path.splitext(file)[1]}".strip()

                new_file_path = os.path.join(dest_folder, new_file_name)
                counter = 1
                while os.path.exists(new_file_path):
                    base, extension = os.path.splitext(new_file_name)
                    new_file_name = f"{base} ({counter}){extension}"
                    new_file_path = os.path.join(dest_folder, new_file_name)
                    counter += 1

                try:
                    shutil.copy2(src_file_path, new_file_path)
                    log_message(f"Renamed and copied {file} to {new_file_name}")
                except Exception as e:
                    log_message(f"ERROR copying file {src_file_path}: {e}")

# Main function
def main():
    src_dir = r"C:\Users\maste\Desktop\CantariBise"
    dest_dir = r"C:\Users\maste\Desktop\CantariBun"
    rename_and_copy_files(src_dir, dest_dir)
    log_message("Process completed.")

if __name__ == "__main__":
    main()
